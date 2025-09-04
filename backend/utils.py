import os
import traceback
import PyPDF2
from openai import OpenAI
import logging

# --- Import models needed for fetching Material content ---
try:
    from database import db, Material # Ensure Material can be imported here
    # If 'db' is not initialized yet or causes circular imports, consider a different approach
    # for fetching Material (e.g., pass app context or use a service function).
    # For now, assume this works within Flask app context.
    DATABASE_ACCESS_AVAILABLE = True
except ImportError:
    DATABASE_ACCESS_AVAILABLE = False
    logging.warning("utils.py: Could not import 'db' or 'Material' from database. Full text from material ID will not work directly here.")


# Configure logging (assuming basicConfig is set in app.py, or set it here)
# logging.basicConfig(level=logging.INFO) # Usually configured once in main app
logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PowerPoint file processing disabled.")

try:
    client = OpenAI() # Uses OPENAI_API_KEY from environment
except Exception as e:
    logger.error(f"Failed to initialize OpenAI client: {e}")
    client = None

# --- Truncation constant ---
# Max characters to send to OpenAI for a single material's full text.
# This is an approximation. A typical token is ~4 chars. 16k model has ~64k char limit.
# GPT-3.5-turbo (4k model) ~16k chars. GPT-3.5-turbo-16k model ~64k chars.
# Let's be conservative for multiple materials and the rest of the prompt.
# For a single material placeholder in a larger prompt:
MAX_CHARS_PER_MATERIAL_CONTEXT = 12000 # Approx 3000 tokens (adjust based on needs and model used)
# For Quiz generation where only one material's text is the main context:
MAX_CHARS_FOR_QUIZ_CONTEXT = 30000 # Approx 7500 tokens (adjust based on needs)

def extract_text(file_path, filename):
    # ... (Keep this function exactly as it was in your working version) ...
    ext = filename.lower().split('.')[-1]; text = ""
    logger.info(f"Attempting to extract text from {filename} (type: {ext})")
    try:
        if ext == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                if reader.is_encrypted: logger.warning(f"Encrypted PDF: {filename}"); return ""
                for page in reader.pages: text += (page.extract_text() or "") + "\n"
        elif ext in ["ppt", "pptx"]:
            if not PPTX_AVAILABLE: logger.warning(f"PPT/PPTX skip: {filename}"); return ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text: text += shape.text + "\n"
        elif ext == "txt":
             with open(file_path, "r", encoding='utf-8', errors='ignore') as f: text = f.read()
        else: logger.warning(f"Unsupported type: {ext}"); return ""
        logger.info(f"Extracted ~{len(text)} chars from {filename}")
        return text.strip()
    except FileNotFoundError: logger.error(f"File not found for extraction: {file_path}"); return ""
    except Exception as e: logger.exception(f"Error extracting text from {filename}: {e}"); return ""


def summarize_text(text, max_length=15000): # Max length of INPUT text to summarize
    # ... (Keep this function exactly as it was) ...
    if not client: logger.error("OpenAI client NI. Cannot summarize."); return "OpenAI client error."
    if not text or not text.strip(): logger.warning("No text to summarize."); return ""
    if len(text) > max_length: logger.warning(f"Text too long ({len(text)}), truncating to {max_length} for summary."); text = text[:max_length]
    try:
        logger.info("Requesting summarization from OpenAI...");
        prompt_message = f"Provide a concise summary (100-150 words) of the following educational material:\n\n{text}\n\nSummary:"
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": "You are an expert summarizer of educational content."}, {"role": "user", "content": prompt_message}],
            temperature=0.3, max_tokens=250, n=1, stop=None,
        )
        summary = response.choices[0].message.content.strip(); logger.info("Summarization successful.")
        return summary
    except Exception as e: logger.exception(f"Error summarizing with OpenAI: {e}"); return "Error during summarization."


def generate_ai_response(system_prompt, user_prompt):
    # ... (Keep this function exactly as it was, ensure it returns content, usage_dict) ...
    if not client: logger.error("OpenAI client NI. Cannot generate."); return "OpenAI client error.", None
    if not user_prompt: return "User prompt required.", None
    try:
        logger.info("Requesting AI response...")
        if not isinstance(system_prompt, str): logger.warning(f"System prompt type {type(system_prompt)}, converting."); system_prompt = str(system_prompt)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo", # Consider gpt-3.5-turbo-0125 for better instruction following if available
            messages=[ {"role": "system", "content": system_prompt or "You are a helpful AI assistant."}, {"role": "user", "content": user_prompt}],
            temperature=0.7, max_tokens=1500, n=1, stop=None, # Increased max_tokens for student chat
        )
        content = response.choices[0].message.content.strip()
        usage = response.usage.model_dump() if response.usage else None
        logger.info(f"AI response OK. Usage: {usage}")
        return content, usage
    except Exception as e: logger.exception(f"Error generating AI response: {e}"); return f"Error: {e}", None


# --- MODIFIED: construct_final_prompt ---
def construct_final_prompt(prompt_structure, user_question):
    """
    Constructs the final system prompt string from the saved structure,
    replacing material placeholders with their full (truncated) text.
    """
    final_prompt_parts = []
    logger.debug(f"Constructing final prompt from structure: {prompt_structure}")

    if not isinstance(prompt_structure, list):
        logger.error(f"Invalid prompt structure type: {type(prompt_structure)}. Expected list.")
        return "Error: Invalid AI assistant configuration.", user_question

    if not DATABASE_ACCESS_AVAILABLE: # Fallback if Material model couldn't be imported
        logger.error("Database access not available in construct_final_prompt. Cannot resolve material placeholders.")
        # Proceed by just joining content, placeholders will remain as text.
        for block in prompt_structure:
            if isinstance(block, dict) and 'content' in block: final_prompt_parts.append(str(block['content']))
            elif isinstance(block, str): final_prompt_parts.append(block)
        system_prompt = "\n\n".join(final_prompt_parts)
        return system_prompt, user_question

    # Iterate through the blocks provided by the teacher
    for block in prompt_structure:
        if isinstance(block, dict):
            block_content = str(block.get('content', ''))
            is_material_block = block.get('isMaterialBlock', False)
            material_id = block.get('materialId', None)

            # Check if it's a material placeholder to be replaced
            # The content should look like: "[USE_FULL_TEXT_FROM_MATERIAL_ID:xxxx-xxxx-xxxx]"
            # Or use the isMaterialBlock flag and materialId field
            if is_material_block and material_id:
                logger.info(f"Found material placeholder for Material ID: {material_id}")
                try:
                    # Fetch the Material object from the database
                    material_obj = db.session.get(Material, material_id) # Use session.get
                    if material_obj and material_obj.extracted_text:
                        logger.debug(f"Material '{material_obj.filename}' found. Using its extracted text.")
                        full_text = material_obj.extracted_text
                        truncated_text = full_text[:MAX_CHARS_PER_MATERIAL_CONTEXT]
                        if len(full_text) > MAX_CHARS_PER_MATERIAL_CONTEXT:
                            logger.warning(f"Material '{material_obj.filename}' text (len {len(full_text)}) was truncated to {MAX_CHARS_PER_MATERIAL_CONTEXT} chars.")
                        # Replace placeholder text or prepend/append material context
                        # For now, let's assume the block_content itself might contain some instruction like "Based on material X:"
                        # So we append the truncated text.
                        # A more robust way would be for the block_content to ONLY be the placeholder ID,
                        # and we construct the text "Context from [Material Name]:\n[Truncated Text]" here.
                        # For current structure where block_content might be "[USE_...]", let's just append.
                        final_prompt_parts.append(f"\n\n--- Context from Material: {material_obj.filename} ---\n{truncated_text}\n--- End Context from Material: {material_obj.filename} ---\n\n")
                        # Also add the original instruction block content if it's not just the placeholder
                        if not block_content.startswith("[USE_FULL_TEXT_FROM_MATERIAL_ID:"):
                            final_prompt_parts.append(block_content)

                    else:
                        logger.warning(f"Material ID {material_id} not found or has no extracted text. Placeholder block content used: '{block_content}'")
                        final_prompt_parts.append(block_content) # Fallback to stored block content
                except Exception as e:
                    logger.exception(f"Error fetching/processing material ID {material_id} for prompt: {e}")
                    final_prompt_parts.append(block_content) # Fallback
            else:
                # It's a regular text block, just add its content
                final_prompt_parts.append(block_content)
        elif isinstance(block, str): # Handle simple string blocks if needed
             final_prompt_parts.append(block)
        else:
            logger.warning(f"Skipping invalid block in prompt structure during final prompt construction: {block}")

    system_prompt = "\n\n".join(filter(None, final_prompt_parts)) # Join non-empty parts
    logger.debug(f"Constructed System Prompt (length {len(system_prompt)}):\n{system_prompt[:500]}...") # Log beginning of prompt
    return system_prompt, user_question